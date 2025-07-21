import streamlit as st
import pandas as pd
from prophet import Prophet
import matplotlib.pyplot as plt
import matplotlib.dates as mdates
from io import BytesIO
import openai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# --- UI Banner ---
st.set_page_config(page_title="Forecast Companion", layout="wide")
st.markdown("## Narrative Companion for Time‑Series Forecasts")
st.markdown('</br>', unsafe_allow_html=True)

@st.cache_data
def fit_prophet(df, date_col, value_col, periods):
    m = Prophet()
    ts = df.rename(columns={date_col: 'ds', value_col: 'y'})[['ds', 'y']]
    m.fit(ts)
    future = m.make_future_dataframe(periods=periods, freq='M')
    forecast = m.predict(future)
    return forecast

@st.cache_data
def generate_insights(prompt, api_key):
    client = openai.OpenAI(api_key=api_key)
    resp = client.chat.completions.create(
        model="gpt-4o",
        messages=[
            {"role": "system", "content": "You are a financial analyst that writes bullet‑point insights with clarity and design in mind."},
            {"role": "user", "content": prompt}
        ],
        temperature=0.7,
        max_tokens=200
    )
    return resp.choices[0].message.content.strip()


def create_ppt(forecast, insights, periods):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    def add_background(slide, color_rgb):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
        )
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*color_rgb)
        shape.line.fill.background()
        slide.shapes._spTree.remove(shape._element)
        slide.shapes._spTree.insert(2, shape._element)

    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    add_background(slide, (245, 245, 245))
    title = slide.shapes.title
    title.text = "Forecast & Narrative Companion"
    tf = title.text_frame.paragraphs[0]
    tf.font.size = Pt(48)
    tf.font.bold = True

    # Chart Slide
    chart_slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_background(chart_slide, (255, 255, 255))
    chart_title = chart_slide.shapes.title
    chart_title.text = "Forecast Plot"
    ct_tf = chart_title.text_frame.paragraphs[0]
    ct_tf.font.size = Pt(36)
    ct_tf.font.color.rgb = RGBColor(0, 70, 122)

    buf = BytesIO()
    plt.figure(figsize=(10,5))
    plt.plot(forecast['ds'], forecast['yhat'], linewidth=2, marker='o', markersize=4)
    start_idx = len(forecast) - periods - 1
    plt.axvline(forecast['ds'][start_idx], color='gray', linestyle='--')
    plt.title("Forecasted Values Over Time", fontsize=16)
    plt.xlabel('Date', fontsize=12)
    plt.ylabel('Value', fontsize=12)
    plt.grid(True, linestyle='--', alpha=0.5)
    locator = mdates.AutoDateLocator()
    formatter = mdates.ConciseDateFormatter(locator)
    ax = plt.gca()
    ax.xaxis.set_major_locator(locator)
    ax.xaxis.set_major_formatter(formatter)
    plt.xticks(rotation=45)
    plt.tight_layout()
    plt.savefig(buf, format='png', transparent=True)
    plt.close()
    buf.seek(0)
    chart_slide.shapes.add_picture(buf, Inches(1), Inches(1.5), width=Inches(11))

    # Insights Slides
    lines = [ln for ln in insights.split("\n") if ln.strip()]
    for idx, line in enumerate(lines, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        add_background(slide, (255, 255, 255))
        title = slide.shapes.title
        title.text = f"Insight {idx}"
        t_tf = title.text_frame.paragraphs[0]
        t_tf.font.size = Pt(32)
        t_tf.font.color.rgb = RGBColor(0, 70, 122)

        body = slide.shapes.placeholders[1].text_frame
        body.clear()
        p = body.add_paragraph()
        p.text = line
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(50, 50, 50)

    out = BytesIO()
    prs.save(out)
    out.seek(0)
    return out

# --- Sidebar ---
st.sidebar.header("🔧 Configuration")
csv = st.sidebar.file_uploader("Upload CSV", type="csv")
if csv:
    df = pd.read_csv(csv, parse_dates=[0])
    date_col = st.sidebar.selectbox("Date column", df.columns, index=0)
    value_col = st.sidebar.selectbox("Value column", df.columns, index=1)
    periods = st.sidebar.slider("Forecast months", 1, 24, 6)
    api_key = st.sidebar.text_input("OpenAI API Key", type="password")

    if st.sidebar.button("Run Forecast"):
        forecast = fit_prophet(df, date_col, value_col, periods)
        fig, ax = plt.subplots(figsize=(12,6), facecolor='none')
        fig.patch.set_alpha(0)
        ax.patch.set_alpha(0)
        ax.plot(forecast['ds'], forecast['yhat'], linewidth=2, marker='o', markersize=4)
        start_idx = len(forecast) - periods - 1
        ax.axvline(forecast['ds'][start_idx], color='gray', linestyle='--', label='Forecast Start')
        ax.set_title("Forecasted Values Over Time", fontsize=12, color='white')
        ax.set_xlabel('Date', fontsize=12, color='white')
        ax.set_ylabel('Value', fontsize=12, color='white')
        ax.grid(True, linestyle='--', alpha=0.5)
        locator = mdates.AutoDateLocator()
        formatter = mdates.ConciseDateFormatter(locator)
        ax.xaxis.set_major_locator(locator)
        ax.xaxis.set_major_formatter(formatter)
        ax.tick_params(axis='x', colors='white')
        ax.tick_params(axis='y', colors='white')
        legend = ax.legend()
        if legend:
            for text in legend.get_texts():
                text.set_color('white')
        fig.autofmt_xdate(rotation=45)
        plt.tight_layout()
        st.pyplot(fig)

        sample = forecast.tail(8)[['ds','yhat']].to_string(index=False)
        prompt = f"Here are dates & values:\n{sample}\nProvide 5 concise bullet-point insights highlighting trends and anomalies."
        insights = generate_insights(prompt, api_key)
        #st.text_area("Insights", insights, height=400)
        st.markdown("### Insights")
        st.markdown(insights)

        pptx_data = create_ppt(forecast, insights, periods)
        st.download_button("Download Presentation", pptx_data, file_name="forecast_narrative.pptx")
else:
    st.info("Upload a CSV to get started.")
